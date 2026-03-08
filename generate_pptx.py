"""
生成 CCAD 项目介绍 PPTX（暖米色 + 珊瑚橙配色方案）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 配色 ──────────────────────────────────────────────────
BG       = RGBColor(0xFA, 0xF9, 0xF6)   # 暖米色背景
CORAL    = RGBColor(0xE5, 0x60, 0x32)   # 珊瑚橙（主色）
CORAL_DK = RGBColor(0xB8, 0x42, 0x1E)   # 深珊瑚（强调）
CORAL_LT = RGBColor(0xFD, 0xEC, 0xE4)   # 浅珊瑚（卡片背景）
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x28, 0x24, 0x21)   # 暖黑（标题）
GRAY     = RGBColor(0x6B, 0x65, 0x60)   # 暖灰（正文）
LGRAY    = RGBColor(0xE2, 0xDB, 0xD3)   # 浅线条
CARD     = RGBColor(0xFF, 0xFF, 0xFF)   # 卡片白

# ── 尺寸 ──────────────────────────────────────────────────
W = Inches(13.333)  # 宽屏 16:9
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=Pt(18), font_color=DARK, bold=False,
             align=PP_ALIGN.LEFT, wrap=True, italic=False,
             font_name="PingFang SC"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return txBox

def add_para(tf, text, font_size=Pt(14), font_color=DARK,
             bold=False, align=PP_ALIGN.LEFT, space_before=Pt(0),
             font_name="PingFang SC"):
    from pptx.util import Pt as P
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    return p

# ══════════════════════════════════════════════════════════════
# Slide 1 – 封面
# ══════════════════════════════════════════════════════════════
def slide_cover(prs):
    s = blank_slide(prs)
    set_bg(s, BG)

    # 左侧大色块
    add_rect(s, Inches(0), Inches(0), Inches(5.8), H, fill_color=CORAL)

    # 左侧装饰线
    add_rect(s, Inches(5.6), Inches(0), Inches(0.04), H, fill_color=CORAL_DK)

    # 右下角大圆（装饰）
    from pptx.util import Inches as I
    circ = s.shapes.add_shape(9, Inches(9.2), Inches(4.5), Inches(5.5), Inches(5.5))  # oval
    circ.fill.solid()
    circ.fill.fore_color.rgb = CORAL_LT
    circ.line.fill.background()

    # 右侧小圆（装饰）
    circ2 = s.shapes.add_shape(9, Inches(5.8), Inches(0.2), Inches(1.5), Inches(1.5))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = CORAL_LT
    circ2.line.fill.background()

    # 左侧 —— 项目标语（白色）
    add_text(s, "智能 CAD\n审图平台",
             Inches(0.55), Inches(1.6), Inches(4.8), Inches(3.2),
             font_size=Pt(52), font_color=WHITE, bold=True,
             align=PP_ALIGN.LEFT, font_name="PingFang SC")

    # 副标题
    add_text(s, "让每一张图纸都经得起推敲",
             Inches(0.55), Inches(4.55), Inches(4.8), Inches(0.6),
             font_size=Pt(18), font_color=RGBColor(0xFF, 0xDF, 0xD0),
             bold=False, align=PP_ALIGN.LEFT)

    # 英文名
    add_text(s, "CCAD  ·  CAD Review System",
             Inches(0.55), Inches(5.4), Inches(5.0), Inches(0.5),
             font_size=Pt(13), font_color=RGBColor(0xFF, 0xC8, 0xB0),
             bold=False, align=PP_ALIGN.LEFT)

    # 右侧说明文字
    add_text(s, "上传 DWG 图纸与目录\nAI 自动比对索引 · 尺寸 · 材料\n生成结构化审图报告",
             Inches(6.2), Inches(2.8), Inches(6.5), Inches(2.0),
             font_size=Pt(19), font_color=GRAY,
             bold=False, align=PP_ALIGN.LEFT, font_name="PingFang SC")

    # 右侧标签
    for i, tag in enumerate(["DWG 解析", "AI 审核", "版本管理", "在线标注"]):
        tx = add_rect(s, Inches(6.2 + i * 1.65), Inches(5.8), Inches(1.45), Inches(0.5),
                      fill_color=CORAL_LT)
        tf = tx.text_frame
        add_text(s, tag, Inches(6.2 + i * 1.65), Inches(5.78), Inches(1.45), Inches(0.52),
                 font_size=Pt(12), font_color=CORAL_DK, bold=True,
                 align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# Slide 2 – 痛点
# ══════════════════════════════════════════════════════════════
def slide_pain(prs):
    s = blank_slide(prs)
    set_bg(s, BG)

    # 顶部装饰条
    add_rect(s, Inches(0), Inches(0), W, Inches(0.08), fill_color=CORAL)

    # 左侧大号数字装饰
    add_text(s, "01", Inches(0.35), Inches(0.4), Inches(2), Inches(1.4),
             font_size=Pt(72), font_color=CORAL_LT, bold=True, align=PP_ALIGN.LEFT)

    # 页标题
    add_text(s, "传统审图，难在哪里？",
             Inches(0.55), Inches(0.38), Inches(8), Inches(0.9),
             font_size=Pt(34), font_color=DARK, bold=True, align=PP_ALIGN.LEFT)

    add_text(s, "工程项目审图的四大核心痛点",
             Inches(0.55), Inches(1.2), Inches(8), Inches(0.5),
             font_size=Pt(15), font_color=GRAY, align=PP_ALIGN.LEFT)

    # 装饰线
    add_rect(s, Inches(0.55), Inches(1.7), Inches(0.5), Inches(0.04), fill_color=CORAL)

    # 4个痛点卡片
    pains = [
        ("耗时费力", "人工逐张比对图纸与目录，\n单个项目动辄数百张，极易遗漏。"),
        ("版本混乱", "多次修版后难以追踪变更，\n新旧版本标注无法直观对比。"),
        ("标准不一", "审图结论依赖个人经验，\n缺乏统一规则与结构化输出。"),
        ("沟通低效", "问题分散于纸质批注，\n无法与图纸位置精确关联。"),
    ]

    for i, (title, desc) in enumerate(pains):
        col = i % 2
        row = i // 2
        cx = Inches(0.55 + col * 6.3)
        cy = Inches(2.0 + row * 2.4)
        cw = Inches(5.9)
        ch = Inches(2.1)

        # 卡片背景
        card = add_rect(s, cx, cy, cw, ch, fill_color=CARD,
                        line_color=LGRAY, line_width=Pt(0.75))

        # 编号色块
        add_rect(s, cx, cy, Inches(0.08), ch, fill_color=CORAL)

        # 序号
        add_text(s, f"0{i+1}", cx + Inches(0.18), cy + Inches(0.18),
                 Inches(0.5), Inches(0.5),
                 font_size=Pt(13), font_color=CORAL, bold=True)

        # 标题
        add_text(s, title,
                 cx + Inches(0.18), cy + Inches(0.55), cw - Inches(0.3), Inches(0.48),
                 font_size=Pt(20), font_color=DARK, bold=True)

        # 描述
        add_text(s, desc,
                 cx + Inches(0.18), cy + Inches(1.0), cw - Inches(0.3), Inches(1.0),
                 font_size=Pt(13.5), font_color=GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════
# Slide 3 – 系统概览
# ══════════════════════════════════════════════════════════════
def slide_overview(prs):
    s = blank_slide(prs)
    set_bg(s, BG)

    add_rect(s, Inches(0), Inches(0), W, Inches(0.08), fill_color=CORAL)

    add_text(s, "02", Inches(0.35), Inches(0.4), Inches(2), Inches(1.4),
             font_size=Pt(72), font_color=CORAL_LT, bold=True)

    add_text(s, "CCAD 是什么？",
             Inches(0.55), Inches(0.38), Inches(9), Inches(0.9),
             font_size=Pt(34), font_color=DARK, bold=True)

    add_text(s, "一站式 CAD 工程图纸智能审图平台",
             Inches(0.55), Inches(1.2), Inches(9), Inches(0.5),
             font_size=Pt(15), font_color=GRAY)

    add_rect(s, Inches(0.55), Inches(1.7), Inches(0.5), Inches(0.04), fill_color=CORAL)

    # 中央流程图：三个大块
    blocks = [
        ("📁", "输入层", "上传 DWG 图纸\n上传 PDF 图册\n上传目录文档", CORAL_LT, CORAL_DK),
        ("⚙️", "处理层", "DWG 解析拆图\n目录三线匹配\nAI 多阶段审核", RGBColor(0xFF, 0xF7, 0xF0), CORAL),
        ("📊", "输出层", "结构化审图报告\n在线画布标注\n多版本对比", RGBColor(0xF5, 0xF0, 0xEB), DARK),
    ]

    for i, (icon, title, desc, bg_c, txt_c) in enumerate(blocks):
        bx = Inches(1.0 + i * 4.1)
        by = Inches(2.2)
        bw = Inches(3.7)
        bh = Inches(4.3)

        add_rect(s, bx, by, bw, bh, fill_color=bg_c,
                 line_color=LGRAY, line_width=Pt(0.75))

        # 顶部色条
        add_rect(s, bx, by, bw, Inches(0.1), fill_color=txt_c)

        add_text(s, icon, bx, by + Inches(0.3), bw, Inches(0.8),
                 font_size=Pt(32), align=PP_ALIGN.CENTER)

        add_text(s, title,
                 bx, by + Inches(1.15), bw, Inches(0.65),
                 font_size=Pt(22), font_color=txt_c, bold=True, align=PP_ALIGN.CENTER)

        add_text(s, desc,
                 bx + Inches(0.2), by + Inches(1.9), bw - Inches(0.4), Inches(2.1),
                 font_size=Pt(14.5), font_color=GRAY, align=PP_ALIGN.CENTER, wrap=True)

        # 箭头（非最后一个）
        if i < 2:
            add_text(s, "→",
                     Inches(4.55 + i * 4.1), by + Inches(1.9), Inches(0.6), Inches(0.6),
                     font_size=Pt(28), font_color=CORAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# Slide 4 – 核心功能
# ══════════════════════════════════════════════════════════════
def slide_features(prs):
    s = blank_slide(prs)
    set_bg(s, BG)

    add_rect(s, Inches(0), Inches(0), W, Inches(0.08), fill_color=CORAL)

    add_text(s, "03", Inches(0.35), Inches(0.4), Inches(2), Inches(1.4),
             font_size=Pt(72), font_color=CORAL_LT, bold=True)

    add_text(s, "四大核心能力",
             Inches(0.55), Inches(0.38), Inches(9), Inches(0.9),
             font_size=Pt(34), font_color=DARK, bold=True)

    add_text(s, "从文件上传到报告生成，全流程覆盖",
             Inches(0.55), Inches(1.2), Inches(9), Inches(0.5),
             font_size=Pt(15), font_color=GRAY)

    add_rect(s, Inches(0.55), Inches(1.7), Inches(0.5), Inches(0.04), fill_color=CORAL)

    features = [
        ("DWG 智能解析", "自动将 DWG 文件按布局拆分为独立图纸，提取图号、标题、尺寸标注等结构化信息，无需人工预处理。", "◈"),
        ("三线自动匹配", "将 目录 × PDF 图纸 × DWG 图纸 三个来源精准关联，智能评分算法处理编号差异与格式不一致问题。", "◉"),
        ("AI 多阶段审核", "依次执行索引断链、尺寸比对、材料表验证，每一步均输出定位信息与问题描述，可追溯可复查。", "◈"),
        ("在线画布标注", "内置笔迹与文字标注工具，支持多审核版本叠加对比，批注与图纸位置精确绑定。", "◉"),
    ]

    for i, (title, desc, icon) in enumerate(features):
        col = i % 2
        row = i // 2
        cx = Inches(0.55 + col * 6.4)
        cy = Inches(2.1 + row * 2.5)

        # 左侧珊瑚大竖线
        add_rect(s, cx, cy, Inches(0.06), Inches(2.1), fill_color=CORAL)

        # 序号
        add_text(s, f"{i+1:02d}",
                 cx + Inches(0.2), cy, Inches(0.7), Inches(0.55),
                 font_size=Pt(26), font_color=CORAL, bold=True)

        add_text(s, title,
                 cx + Inches(0.2), cy + Inches(0.5), Inches(5.6), Inches(0.55),
                 font_size=Pt(19), font_color=DARK, bold=True)

        add_text(s, desc,
                 cx + Inches(0.2), cy + Inches(1.1), Inches(5.8), Inches(1.0),
                 font_size=Pt(13), font_color=GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════
# Slide 5 – 审核流程
# ══════════════════════════════════════════════════════════════
def slide_workflow(prs):
    s = blank_slide(prs)
    set_bg(s, BG)

    add_rect(s, Inches(0), Inches(0), W, Inches(0.08), fill_color=CORAL)

    add_text(s, "04", Inches(0.35), Inches(0.4), Inches(2), Inches(1.4),
             font_size=Pt(72), font_color=CORAL_LT, bold=True)

    add_text(s, "7 步 AI 审核流程",
             Inches(0.55), Inches(0.38), Inches(9), Inches(0.9),
             font_size=Pt(34), font_color=DARK, bold=True)

    add_text(s, "系统化拆解审图任务，每一环节实时可见",
             Inches(0.55), Inches(1.2), Inches(9), Inches(0.5),
             font_size=Pt(15), font_color=GRAY)

    add_rect(s, Inches(0.55), Inches(1.7), Inches(0.5), Inches(0.04), fill_color=CORAL)

    steps = [
        ("准备检查", "校验三线匹配完整性"),
        ("提取图纸", "构建每张图的上下文"),
        ("规划路径", "生成有向任务审核图"),
        ("索引核对", "断链与错编检查"),
        ("尺寸比对", "跨版本尺寸差异识别"),
        ("材料验证", "材料表与图纸一致性"),
        ("生成报告", "结构化输出·可导出"),
    ]

    # 连接横线
    add_rect(s, Inches(0.85), Inches(4.02), Inches(11.8), Inches(0.06), fill_color=LGRAY)

    for i, (title, desc) in enumerate(steps):
        sx = Inches(0.55 + i * 1.75)
        sy = Inches(2.1)
        sw = Inches(1.55)

        is_active = i in [3, 4, 5]  # 核心审核步骤高亮
        bg_c = CORAL if is_active else CARD
        txt_c = WHITE if is_active else DARK

        # 节点圆形
        circ = s.shapes.add_shape(9,
            sx + Inches(0.375), Inches(3.72),
            Inches(0.8), Inches(0.8))
        circ.fill.solid()
        circ.fill.fore_color.rgb = CORAL if is_active else LGRAY
        circ.line.fill.background()

        # 步骤编号
        add_text(s, str(i + 1),
                 sx + Inches(0.375), Inches(3.74),
                 Inches(0.8), Inches(0.75),
                 font_size=Pt(16), font_color=WHITE if is_active else GRAY,
                 bold=True, align=PP_ALIGN.CENTER)

        # 步骤卡片
        card = add_rect(s, sx, sy, sw, Inches(1.45),
                        fill_color=bg_c,
                        line_color=CORAL if is_active else LGRAY,
                        line_width=Pt(1.0 if is_active else 0.75))

        add_text(s, title,
                 sx, sy + Inches(0.2), sw, Inches(0.65),
                 font_size=Pt(14.5), font_color=txt_c, bold=True,
                 align=PP_ALIGN.CENTER)

        add_text(s, desc,
                 sx + Inches(0.05), sy + Inches(0.75), sw - Inches(0.1), Inches(0.65),
                 font_size=Pt(10.5), font_color=WHITE if is_active else GRAY,
                 align=PP_ALIGN.CENTER, wrap=True)

        # 描述文字（下方）
        add_text(s, desc,
                 sx, Inches(4.72), sw, Inches(0.8),
                 font_size=Pt(11), font_color=GRAY,
                 align=PP_ALIGN.CENTER, wrap=True)

    # 注释
    add_text(s, "★ 高亮步骤为 AI 深度比对阶段，支持单步重跑",
             Inches(0.55), Inches(6.7), Inches(8), Inches(0.5),
             font_size=Pt(12), font_color=CORAL, italic=True)


# ══════════════════════════════════════════════════════════════
# Slide 6 – 技术亮点（版本管理 & 标注）
# ══════════════════════════════════════════════════════════════
def slide_highlights(prs):
    s = blank_slide(prs)
    set_bg(s, BG)

    add_rect(s, Inches(0), Inches(0), W, Inches(0.08), fill_color=CORAL)

    add_text(s, "05", Inches(0.35), Inches(0.4), Inches(2), Inches(1.4),
             font_size=Pt(72), font_color=CORAL_LT, bold=True)

    add_text(s, "设计亮点",
             Inches(0.55), Inches(0.38), Inches(9), Inches(0.9),
             font_size=Pt(34), font_color=DARK, bold=True)

    add_rect(s, Inches(0.55), Inches(1.7), Inches(0.5), Inches(0.04), fill_color=CORAL)

    highlights = [
        {
            "title": "多版本叠加对比",
            "points": [
                "每次审核自动创建独立版本快照",
                "在同一画布上叠加显示不同版本标注",
                "支持 2 个版本快捷切换 / 超过 2 个版本下拉选择",
                "版本间一键对比，变更清晰可见",
            ]
        },
        {
            "title": "画布标注工具",
            "points": [
                "内置画笔工具，支持调节粗细（2–20px）",
                "文字标注，字号可调（12–36px）",
                "支持撤销与清空，操作轻量高效",
                "标注随图纸坐标绑定，缩放不失位",
            ]
        },
        {
            "title": "三线匹配引擎",
            "points": [
                "目录 × PDF × DWG 三端来源自动关联",
                "智能文本归一化处理编号格式差异",
                "SequenceMatcher 模糊评分 + 规则加权",
                "可视化匹配表格，支持人工修正",
            ]
        },
    ]

    for i, h in enumerate(highlights):
        cx = Inches(0.55 + i * 4.3)
        cy = Inches(2.0)
        cw = Inches(4.0)
        ch = Inches(4.9)

        add_rect(s, cx, cy, cw, ch, fill_color=CARD,
                 line_color=LGRAY, line_width=Pt(0.75))

        # 顶部色块
        add_rect(s, cx, cy, cw, Inches(0.62), fill_color=CORAL_LT)
        add_rect(s, cx, cy, Inches(0.08), Inches(0.62), fill_color=CORAL)

        add_text(s, h["title"],
                 cx + Inches(0.2), cy + Inches(0.1), cw - Inches(0.25), Inches(0.5),
                 font_size=Pt(17), font_color=CORAL_DK, bold=True)

        for j, point in enumerate(h["points"]):
            py = cy + Inches(0.85 + j * 0.96)
            # 小圆点
            dot = s.shapes.add_shape(9,
                cx + Inches(0.2), py + Inches(0.12),
                Inches(0.12), Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = CORAL
            dot.line.fill.background()

            add_text(s, point,
                     cx + Inches(0.42), py, cw - Inches(0.55), Inches(0.85),
                     font_size=Pt(13), font_color=GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════
# Slide 7 – 技术架构
# ══════════════════════════════════════════════════════════════
def slide_tech(prs):
    s = blank_slide(prs)
    set_bg(s, BG)

    add_rect(s, Inches(0), Inches(0), W, Inches(0.08), fill_color=CORAL)

    add_text(s, "06", Inches(0.35), Inches(0.4), Inches(2), Inches(1.4),
             font_size=Pt(72), font_color=CORAL_LT, bold=True)

    add_text(s, "技术架构",
             Inches(0.55), Inches(0.38), Inches(9), Inches(0.9),
             font_size=Pt(34), font_color=DARK, bold=True)

    add_text(s, "现代化全栈工程，前后端解耦，模块清晰",
             Inches(0.55), Inches(1.2), Inches(9), Inches(0.5),
             font_size=Pt(15), font_color=GRAY)

    add_rect(s, Inches(0.55), Inches(1.7), Inches(0.5), Inches(0.04), fill_color=CORAL)

    layers = [
        ("前端", "React 18 + TypeScript · Vite · Tailwind CSS · Radix UI · Framer Motion", CORAL_LT, CORAL_DK),
        ("后端", "FastAPI (Python) · SQLAlchemy ORM · Pydantic v2 · Async 异步架构", RGBColor(0xF5, 0xF0, 0xEB), DARK),
        ("文件处理", "ezdxf（DWG/DXF 解析）· PyMuPDF（PDF 渲染）· OpenPyXL（Excel 目录）", RGBColor(0xFF, 0xF7, 0xF0), CORAL),
        ("AI 审核", "OpenAI / 兼容大模型 API · 有向任务图调度 · 结构化 JSON 输出", CORAL_LT, CORAL_DK),
        ("存储", "SQLite / PostgreSQL（项目 & 目录）· 文件系统（图纸 & 标注 JSON）", RGBColor(0xF5, 0xF0, 0xEB), DARK),
    ]

    for i, (layer, desc, bg_c, txt_c) in enumerate(layers):
        ly = Inches(2.1 + i * 1.02)

        add_rect(s, Inches(0.55), ly, Inches(2.1), Inches(0.82), fill_color=txt_c if i % 2 == 0 else CORAL)
        add_text(s, layer,
                 Inches(0.55), ly + Inches(0.15), Inches(2.1), Inches(0.55),
                 font_size=Pt(15), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        add_rect(s, Inches(2.65), ly, Inches(10.1), Inches(0.82),
                 fill_color=bg_c, line_color=LGRAY, line_width=Pt(0.5))

        add_text(s, desc,
                 Inches(2.85), ly + Inches(0.15), Inches(9.8), Inches(0.55),
                 font_size=Pt(13.5), font_color=GRAY)


# ══════════════════════════════════════════════════════════════
# Slide 8 – 结语
# ══════════════════════════════════════════════════════════════
def slide_closing(prs):
    s = blank_slide(prs)
    set_bg(s, BG)

    # 右侧大色块
    add_rect(s, Inches(7.5), Inches(0), Inches(5.833), H, fill_color=CORAL)

    # 装饰线
    add_rect(s, Inches(7.46), Inches(0), Inches(0.04), H, fill_color=CORAL_DK)

    # 左侧装饰圆
    circ = s.shapes.add_shape(9, Inches(-1.5), Inches(-1.5), Inches(5), Inches(5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = CORAL_LT
    circ.line.fill.background()

    # 左侧文字
    add_text(s, "让审图\n回归本质",
             Inches(0.6), Inches(1.5), Inches(6.5), Inches(2.5),
             font_size=Pt(48), font_color=DARK, bold=True, align=PP_ALIGN.LEFT)

    add_text(s, "CCAD 不是要取代工程师的判断，\n而是让工程师从繁琐的比对中解放出来，\n专注于真正需要专业经验的决策。",
             Inches(0.6), Inches(4.1), Inches(6.5), Inches(2.0),
             font_size=Pt(15.5), font_color=GRAY, align=PP_ALIGN.LEFT, wrap=True)

    # 右侧关键指标
    stats = [
        ("7", "审核阶段全覆盖"),
        ("3", "文件来源三线匹配"),
        ("∞", "多版本历史回溯"),
    ]

    for i, (num, label) in enumerate(stats):
        sy = Inches(1.8 + i * 1.8)

        add_text(s, num,
                 Inches(8.0), sy, Inches(2.5), Inches(1.1),
                 font_size=Pt(60), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        add_text(s, label,
                 Inches(8.0), sy + Inches(1.0), Inches(5.0), Inches(0.5),
                 font_size=Pt(14), font_color=RGBColor(0xFF, 0xDF, 0xD0),
                 align=PP_ALIGN.CENTER)

        if i < 2:
            add_rect(s, Inches(9.0), sy + Inches(1.55), Inches(3.0), Inches(0.03),
                     fill_color=RGBColor(0xFF, 0xAA, 0x88))

    # 右下角
    add_text(s, "CCAD  ·  CAD Review System",
             Inches(7.7), Inches(6.8), Inches(5.3), Inches(0.45),
             font_size=Pt(12), font_color=RGBColor(0xFF, 0xC8, 0xB0),
             align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# 主程序
# ══════════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    slide_cover(prs)
    slide_pain(prs)
    slide_overview(prs)
    slide_features(prs)
    slide_workflow(prs)
    slide_highlights(prs)
    slide_tech(prs)
    slide_closing(prs)

    out = "/Users/harry/@dev/ccad/CCAD项目介绍.pptx"
    prs.save(out)
    print(f"✓ 已生成：{out}")
    print(f"  共 {len(prs.slides)} 张幻灯片")

if __name__ == "__main__":
    main()
